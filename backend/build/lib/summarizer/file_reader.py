# summarizer/file_reader.py

import docx
import csv
import PyPDF2
from pptx import Presentation
from logger.logger import get_logger

logger = get_logger(__name__)

def read_docx(file_path):
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        logger.error(f"Failed to read DOCX file at {file_path}: {str(e)}")
        raise

def read_pdf(file_path):
    try:
        pdf_text = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                pdf_text.append(page.extract_text())
        return '\n'.join(pdf_text)
    except Exception as e:
        logger.error(f"Failed to read PDF file at {file_path}: {str(e)}")
        raise

def read_ppt(file_path):
    try:
        ppt_text = []
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text.append(shape.text)
        return '\n'.join(ppt_text)
    except Exception as e:
        logger.error(f"Failed to read PPT file at {file_path}: {str(e)}")
        raise

def read_csv(file_path):
    try:
        csv_text = []
        with open(file_path, mode='r', newline='', encoding='utf-8') as file:
            reader = csv.reader(file)
            for row in reader:
                csv_text.append(' '.join(row))
        return '\n'.join(csv_text)
    except Exception as e:
        logger.error(f"Failed to read CSV file at {file_path}: {str(e)}")
        raise

def read_file(file_path):
    if file_path.endswith('.docx'):
        return read_docx(file_path)
    elif file_path.endswith('.pdf'):
        return read_pdf(file_path)
    elif file_path.endswith('.pptx'):
        return read_ppt(file_path)
    elif file_path.endswith('.csv'):
        return read_csv(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")
